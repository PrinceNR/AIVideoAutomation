from dataclasses import dataclass
from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

from presentation.animations.visual_animation_settings import (
    VisualAnimationSettings,
)


@dataclass(frozen=True)
class VisualAnimationSpec:
    shape_name: str
    semantic_element: str
    effect_id: int
    duration: float
    direction: int | None = None
    text_unit_effect: int | None = None
    required: bool = False


@dataclass(frozen=True)
class SlideTransitionSpec:
    entry_effect: int
    speed: int


class AnimationTemplateError(RuntimeError):
    """Raised when required animation semantics are missing or ambiguous."""


class VisualAnimationPlanner:

    SLIDES_PER_WORD = 4

    EXCLUDED_SHAPE_NAMES = {
        "PROGRESS_TRACK",
        "PROGRESS_FILL",
    }

    INTRO_PLACEHOLDER_RULES = (
        (
            "word",
            ("{{WORD}}",),
        ),
        (
            "pronunciation",
            ("{{PRONUNCIATION}}",),
        ),
        (
            "meaning",
            ("{{MEANING}}",),
        ),
        (
            "translation",
            (
                "{{HINDI}}",
                "{{MALAYALAM}}",
                "{{TAMIL}}",
            ),
        ),
        (
            "verb_form",
            (
                "{{BASE_FORM}}",
                "{{PRESENT_FORM}}",
                "{{PAST_FORM}}",
            ),
        ),
    )

    INTRO_SETTING_FIELDS = {
        "word": (
            "word_effect",
            "word_duration",
            "word_direction",
        ),
        "pronunciation": (
            "pronunciation_effect",
            "pronunciation_duration",
            None,
        ),
        "meaning": (
            "meaning_effect",
            "meaning_duration",
            None,
        ),
        "translation": (
            "translation_effect",
            "translation_duration",
            None,
        ),
        "verb_form": (
            "verb_form_effect",
            "verb_form_duration",
            None,
        ),
    }

    SENTENCE_SHAPES = {
        2: (
            "PAST_SENTENCE",
            "past_sentence",
        ),
        3: (
            "PRESENT_SENTENCE",
            "present_sentence",
        ),
        4: (
            "FUTURE_SENTENCE",
            "future_sentence",
        ),
    }

    def __init__(
        self,
        settings=None,
    ):
        self.settings = (
            settings
            or VisualAnimationSettings.from_project_config()
        )

    def transition_spec(
        self,
        slide_within_word
    ):
        return SlideTransitionSpec(
            entry_effect=self.settings.transition_id(
                self.settings.new_word_transition
                if slide_within_word == 1
                else self.settings.continuation_transition
            ),
            speed=(
                self.settings.transition_speed_id()
            ),
        )

    def slide_within_word(
        self,
        slide_index
    ):
        return (
            (slide_index - 1)
            % self.SLIDES_PER_WORD
        ) + 1

    def build_template_plan(
        self,
        template_path
    ):
        presentation = Presentation(
            str(Path(template_path))
        )

        slides = list(
            presentation.slides
        )

        if len(slides) != self.SLIDES_PER_WORD:
            raise AnimationTemplateError(
                "Presentation animation template must contain "
                f"exactly {self.SLIDES_PER_WORD} slides; "
                f"found {len(slides)}."
            )

        return [
            self.plan_slide(
                slide,
                slide_within_word,
            )
            for slide_within_word, slide in enumerate(
                slides,
                start=1,
            )
        ]

    def plan_slide(
        self,
        slide,
        slide_within_word,
    ):
        shapes = list(
            self._iter_shapes(
                slide.shapes
            )
        )

        if slide_within_word == 1:
            return self._plan_intro_slide(
                shapes
            )

        return [
            self._required_sentence_spec(
                shapes,
                slide_within_word,
            )
        ]

    def _plan_intro_slide(
        self,
        shapes
    ):
        specs = []

        for shape in shapes:
            text = (
                shape.text
                if getattr(
                    shape,
                    "has_text_frame",
                    False
                )
                else ""
            )

            spec = self.plan_intro_shape(
                shape.name,
                text
            )

            if spec is not None:
                specs.append(spec)

        return specs

    def plan_intro_shape(
        self,
        shape_name,
        source_text=""
    ):
        if shape_name in self.EXCLUDED_SHAPE_NAMES:
            return None

        if shape_name == "VOCAB_IMAGE":
            return VisualAnimationSpec(
                shape_name=shape_name,
                semantic_element="image",
                effect_id=self.settings.effect_id(
                    self.settings.image_effect
                ),
                duration=self.settings.image_duration,
            )

        normalized_text = "".join(
            str(source_text).split()
        ).upper()

        for (
            semantic_element,
            placeholders,
        ) in self.INTRO_PLACEHOLDER_RULES:
            if any(
                placeholder in normalized_text
                for placeholder in placeholders
            ):
                (
                    effect_field,
                    duration_field,
                    direction_field,
                ) = self.INTRO_SETTING_FIELDS[
                    semantic_element
                ]
                effect_name = getattr(
                    self.settings,
                    effect_field,
                )
                direction_name = (
                    getattr(
                        self.settings,
                        direction_field,
                    )
                    if direction_field is not None
                    else None
                )

                return VisualAnimationSpec(
                    shape_name=shape_name,
                    semantic_element=semantic_element,
                    effect_id=self.settings.effect_id(
                        effect_name
                    ),
                    duration=getattr(
                        self.settings,
                        duration_field,
                    ),
                    direction=self.settings.direction_id(
                        direction_name
                    ),
                )

        return None

    def _required_sentence_spec(
        self,
        shapes,
        slide_within_word,
    ):
        try:
            shape_name, semantic_element = (
                self.SENTENCE_SHAPES[
                    slide_within_word
                ]
            )
        except KeyError as error:
            raise AnimationTemplateError(
                "Unsupported slide-within-word position for "
                f"animation planning: {slide_within_word}."
            ) from error

        matches = [
            shape
            for shape in shapes
            if shape.name == shape_name
        ]

        if not matches:
            raise AnimationTemplateError(
                "Presentation animation template slide "
                f"{slide_within_word} is missing required "
                f"semantic shape '{shape_name}'."
            )

        if len(matches) > 1:
            raise AnimationTemplateError(
                "Presentation animation template slide "
                f"{slide_within_word} must contain exactly one "
                f"semantic shape '{shape_name}'; "
                f"found {len(matches)}."
            )

        return VisualAnimationSpec(
            shape_name=shape_name,
            semantic_element=semantic_element,
            effect_id=self.settings.effect_id(
                self.settings.sentence_effect
            ),
            duration=self.settings.sentence_duration,
            direction=self.settings.direction_id(
                self.settings.sentence_direction
            ),
            text_unit_effect=self.settings.text_unit_effect(),
            required=True,
        )

    def _iter_shapes(
        self,
        shapes
    ):
        for shape in shapes:
            if (
                shape.shape_type
                == MSO_SHAPE_TYPE.GROUP
            ):
                yield from self._iter_shapes(
                    shape.shapes
                )
            else:
                yield shape
