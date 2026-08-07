from pathlib import Path

from pptx import Presentation

from presentation.embedders.audio_embedder import AudioEmbedder


presentation = Presentation()

slide = presentation.slides.add_slide(
    presentation.slide_layouts[6]   # Blank slide
)

audio_embedder = AudioEmbedder()

audio_embedder.embed(
    slide,
    Path("output/college/audio/campus/pronunciation.mp3"),
    0.0
)

presentation.save(
    "research/python_audio.pptx"
)

print("Research presentation created.")