from pptx.enum.shapes import MSO_SHAPE_TYPE


class ProgressTemplateError(RuntimeError):
    """Raised when the presentation template lacks a progress shape."""


class ProgressProcessor:

    TRACK_SHAPE_NAME = "PROGRESS_TRACK"
    FILL_SHAPE_NAME = "PROGRESS_FILL"

    def process(
        self,
        slide,
        slide_definition,
        word,
        word_number,
        total_words,
        timeline
    ):
        track = self._find_shape(
            slide,
            self.TRACK_SHAPE_NAME
        )

        if track is None:
            raise ProgressTemplateError(
                "Presentation template is missing required "
                "semantic shape 'PROGRESS_TRACK'."
            )

        fill = self._find_shape(
            slide,
            self.FILL_SHAPE_NAME
        )

        if fill is None:
            raise ProgressTemplateError(
                "Presentation template is missing required "
                "semantic shape 'PROGRESS_FILL'."
            )

        progress = self._calculate_progress(
            word_number,
            total_words
        )

        # Width is the only fill-shape property changed.
        fill.width = round(
            track.width * progress
        )

    def _calculate_progress(
        self,
        word_number,
        total_words
    ):
        if total_words <= 0:
            return 0.0

        progress = word_number / total_words

        return max(
            0.0,
            min(1.0, progress)
        )

    def _find_shape(
        self,
        slide,
        semantic_name
    ):
        for shape in self._iter_shapes(
            slide.shapes
        ):
            if shape.name == semantic_name:
                return shape

        return None

    def _iter_shapes(
        self,
        shapes
    ):
        for shape in shapes:
            if (
                shape.shape_type
                == MSO_SHAPE_TYPE.GROUP
            ):
                yield from self._iter_shapes(
                    shape.shapes
                )
            else:
                yield shape
