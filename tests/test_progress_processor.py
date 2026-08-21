import unittest
from types import SimpleNamespace

from pptx.enum.shapes import MSO_SHAPE_TYPE

from presentation.processors.progress_processor import (
    ProgressProcessor,
    ProgressTemplateError,
)
from presentation.slide_group_renderer import SlideGroupRenderer
from presentation.template_definition_loader import (
    TemplateDefinitionLoader,
)


class FakeShape:

    def __init__(
        self,
        name,
        width,
        left=0,
        top=0,
        height=0,
        shape_type=MSO_SHAPE_TYPE.AUTO_SHAPE,
        shapes=None,
    ):
        self.name = name
        self.width = width
        self.left = left
        self.top = top
        self.height = height
        self.shape_type = shape_type
        self.shapes = shapes or []
        self.fill = object()
        self.line = object()
        self.transparency = object()
        self.effects = object()


class FakeSlide:

    def __init__(self, shapes):
        self.shapes = shapes


class ProgressProcessorTests(unittest.TestCase):

    TRACK_WIDTH = 1000

    def setUp(self):
        self.processor = ProgressProcessor()

    def _make_slide(
        self,
        include_track=True,
        include_fill=True,
    ):
        shapes = []

        if include_track:
            shapes.append(
                FakeShape(
                    "PROGRESS_TRACK",
                    self.TRACK_WIDTH,
                )
            )

        fill = None

        if include_fill:
            fill = FakeShape(
                "PROGRESS_FILL",
                self.TRACK_WIDTH,
                left=25,
                top=50,
                height=75,
            )
            shapes.append(fill)

        return FakeSlide(shapes), fill

    def _process(self, word_number, total_words):
        slide, fill = self._make_slide()

        self.processor.process(
            slide=slide,
            slide_definition=None,
            word=None,
            word_number=word_number,
            total_words=total_words,
            timeline=None,
        )

        return fill

    def test_word_1_of_10_is_ten_percent(self):
        fill = self._process(1, 10)

        self.assertEqual(fill.width, 100)

    def test_word_5_of_10_is_fifty_percent(self):
        fill = self._process(5, 10)

        self.assertEqual(fill.width, 500)

    def test_word_10_of_10_fills_complete_track(self):
        fill = self._process(10, 10)

        self.assertEqual(fill.width, self.TRACK_WIDTH)

    def test_arbitrary_lesson_sizes_are_supported(self):
        for word_number, total_words in (
            (1, 3),
            (2, 3),
            (4, 7),
            (7, 7),
        ):
            with self.subTest(
                word_number=word_number,
                total_words=total_words,
            ):
                fill = self._process(
                    word_number,
                    total_words,
                )
                expected = round(
                    self.TRACK_WIDTH
                    * word_number
                    / total_words
                )

                self.assertEqual(fill.width, expected)

    def test_all_four_slides_for_word_have_identical_progress(self):
        slides_and_fills = [
            self._make_slide()
            for _ in range(4)
        ]
        slides = [
            slide
            for slide, _ in slides_and_fills
        ]
        definitions = [
            SimpleNamespace(
                type="TEST",
                processors=["progress"],
                audio=None,
            )
            for _ in range(4)
        ]

        SlideGroupRenderer().render(
            slides=slides,
            slide_definitions=definitions,
            word=None,
            word_number=5,
            total_words=10,
        )

        self.assertEqual(
            [
                fill.width
                for _, fill in slides_and_fills
            ],
            [500, 500, 500, 500],
        )

    def test_fill_position_height_and_appearance_remain_unchanged(self):
        slide, fill = self._make_slide()
        original_properties = (
            fill.left,
            fill.top,
            fill.height,
            fill.fill,
            fill.line,
            fill.transparency,
            fill.effects,
        )

        self.processor.process(
            slide,
            None,
            None,
            5,
            10,
            None,
        )

        self.assertEqual(
            (
                fill.left,
                fill.top,
                fill.height,
                fill.fill,
                fill.line,
                fill.transparency,
                fill.effects,
            ),
            original_properties,
        )

    def test_grouped_progress_shapes_are_found_recursively(self):
        track = FakeShape(
            "PROGRESS_TRACK",
            self.TRACK_WIDTH,
        )
        fill = FakeShape(
            "PROGRESS_FILL",
            self.TRACK_WIDTH,
        )
        group = FakeShape(
            "Progress Group",
            self.TRACK_WIDTH,
            shape_type=MSO_SHAPE_TYPE.GROUP,
            shapes=[track, fill],
        )

        self.processor.process(
            FakeSlide([group]),
            None,
            None,
            1,
            10,
            None,
        )

        self.assertEqual(fill.width, 100)

    def test_missing_progress_track_raises_template_error(self):
        slide, _ = self._make_slide(
            include_track=False
        )

        with self.assertRaisesRegex(
            ProgressTemplateError,
            "PROGRESS_TRACK",
        ):
            self.processor.process(
                slide,
                None,
                None,
                1,
                10,
                None,
            )

    def test_missing_progress_fill_raises_template_error(self):
        slide, _ = self._make_slide(
            include_fill=False
        )

        with self.assertRaisesRegex(
            ProgressTemplateError,
            "PROGRESS_FILL",
        ):
            self.processor.process(
                slide,
                None,
                None,
                1,
                10,
                None,
            )

    def test_all_template_slides_enable_progress_processor(self):
        definition = TemplateDefinitionLoader().load(
            "templates/vocabulary/template_definition.json"
        )

        self.assertEqual(len(definition.slides), 4)
        self.assertTrue(
            all(
                "progress" in slide.processors
                for slide in definition.slides
            )
        )


if __name__ == "__main__":
    unittest.main()
