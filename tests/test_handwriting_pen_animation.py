import io
import tempfile
import unittest
from contextlib import redirect_stdout
from dataclasses import replace
from pathlib import Path
from types import SimpleNamespace
from unittest.mock import patch

from PIL import Image
from pptx.enum.shapes import MSO_SHAPE_TYPE

import config as project_config

from presentation.animations.handwriting_pen_asset import (
    HandwritingPenAssetNormalizer,
)
from presentation.animations.visual_animation_planner import (
    AnimationTemplateError,
    VisualAnimationPlanner,
)
from presentation.animations.visual_animation_presentation_processor import (
    VisualAnimationPresentationProcessor,
)
from presentation.animations.visual_animation_settings import (
    VisualAnimationSettings,
)
from presentation.patchers.sentence_letter_timing_patcher import (
    SentenceLetterTimingPatcher,
)
from presentation.video_presentation_processor import (
    VideoPresentationProcessor,
)


class FakeTiming:

    def __init__(self, trigger_type=3, delay=0.0, duration=0.0):
        self.TriggerType = trigger_type
        self.TriggerDelayTime = delay
        self.Duration = duration


class FakeMotionEffect:

    def __init__(self):
        self.ByX = None
        self.ByY = None


class FakeBehavior:

    def __init__(self, behavior_type):
        self.Type = behavior_type
        self.MotionEffect = FakeMotionEffect()


class FakeBehaviors:

    def __init__(self):
        self.items = []

    def Add(self, behavior_type):
        behavior = FakeBehavior(behavior_type)
        self.items.append(behavior)
        return behavior


class FakeEffect:

    def __init__(
        self,
        shape=None,
        effect_type=None,
        trigger_type=3,
        delay=0.0,
        duration=0.0,
        sequence=None,
    ):
        self.Shape = shape
        self.EffectType = effect_type
        self.Timing = FakeTiming(
            trigger_type,
            delay,
            duration,
        )
        self.Behaviors = FakeBehaviors()
        self.EffectInformation = SimpleNamespace(
            TextUnitEffect=0,
            TextLevelEffect=0,
        )
        self.Exit = 0
        self.sequence = sequence

    def Delete(self):
        self.sequence.effects.remove(self)

    def MoveTo(self, position):
        self.sequence.effects.remove(self)
        self.sequence.effects.insert(position - 1, self)


class FakeSequence:

    def __init__(self, effects=None):
        self.effects = list(effects or [])

        for effect in self.effects:
            effect.sequence = self

    @property
    def Count(self):
        return len(self.effects)

    def Item(self, index):
        return self.effects[index - 1]

    def AddEffect(
        self,
        shape,
        effect_id,
        level,
        trigger,
        index=None,
    ):
        effect = FakeEffect(
            shape=shape,
            effect_type=effect_id,
            trigger_type=trigger,
            sequence=self,
        )

        if index is None:
            self.effects.append(effect)
        else:
            self.effects.insert(index - 1, effect)

        return effect

    def ConvertToTextUnitEffect(self, effect, text_unit):
        effect.EffectInformation.TextUnitEffect = text_unit
        settings = effect.Shape.AnimationSettings
        settings.Animate = -1
        # Modern PowerPoint exposes its actual Fade through the timeline;
        # this legacy projection remains an Appear-family value.
        settings.EntryEffect = 257
        settings.TextLevelEffect = 16
        settings.TextUnitEffect = 2
        return effect


class FakeCollection:

    def __init__(self, items=None):
        self.items = list(items or [])

    @property
    def Count(self):
        return len(self.items)

    def Item(self, index):
        return self.items[index - 1]


class FakeFill:

    def __init__(self, rgb=0xFFFFFF):
        self.Type = 1
        self.ForeColor = SimpleNamespace(RGB=rgb)
        self.Transparency = 0.0

    def Solid(self):
        self.Type = 1


class FakeTextLine:

    def __init__(self, left, top, width, height, text="Text"):
        self.BoundLeft = left
        self.BoundTop = top
        self.BoundWidth = width
        self.BoundHeight = height
        self.Text = text


class FakeTextLines:

    def __init__(self, lines):
        self.lines = list(lines)

    @property
    def Count(self):
        return len(self.lines)

    def __call__(self, line_index, length=1):
        if length != 1:
            raise ValueError("Fake text lines support one line at a time")

        return self.lines[line_index - 1]

    def Item(self, line_index):
        return self.lines[line_index - 1]


class FakeTextRange(FakeTextLine):

    def __init__(self, lines):
        lines = list(lines)
        left = min(line.BoundLeft for line in lines)
        top = min(line.BoundTop for line in lines)
        right = max(
            line.BoundLeft + line.BoundWidth
            for line in lines
        )
        bottom = max(
            line.BoundTop + line.BoundHeight
            for line in lines
        )
        super().__init__(
            left,
            top,
            right - left,
            bottom - top,
            "\r".join(line.Text for line in lines),
        )
        self.Lines = FakeTextLines(lines)


class FakeShape:

    def __init__(
        self,
        name,
        left=100.0,
        top=200.0,
        width=480.0,
        height=30.0,
        line_count=1,
        text="A single rendered sentence.",
        line_bounds=None,
        line_texts=None,
        shape_id=1,
    ):
        self.Name = name
        self.Id = shape_id
        self.Type = 1
        self.GroupItems = FakeCollection()
        self.Left = left
        self.Top = top
        self.Width = width
        self.Height = height
        if line_bounds is None:
            rendered_width = min(width * 0.65, 300.0)
            line_bounds = [
                (
                    left + 8.0,
                    top + 4.0 + index * 24.0,
                    rendered_width,
                    20.0,
                )
                for index in range(line_count)
            ]

        if line_texts is None:
            line_texts = [text]

            if len(line_bounds) > 1:
                line_texts.extend(
                    f"line {index}"
                    for index in range(2, len(line_bounds) + 1)
                )

        if len(line_texts) != len(line_bounds):
            raise ValueError(
                "line_texts must match the rendered line count"
            )

        rendered_lines = [
            FakeTextLine(*bounds, line_text)
            for bounds, line_text in zip(line_bounds, line_texts)
        ]
        self.TextFrame2 = SimpleNamespace(
            TextRange=FakeTextRange(rendered_lines)
        )
        self.TextFrame = SimpleNamespace(
            TextRange=SimpleNamespace(Text=text)
        )
        self.AnimationSettings = SimpleNamespace(
            Animate=0,
            EntryEffect=0,
            TextLevelEffect=0,
            TextUnitEffect=0,
        )
        self.Fill = FakeFill()
        self.Line = SimpleNamespace(Visible=-1)
        self.LockAspectRatio = 0
        self.ZOrderPosition = 0
        self.owner = None
        self.z_order_calls = []

    def ZOrder(self, command):
        self.z_order_calls.append(command)
        self.ZOrderPosition = len(self.owner.items) + 1

    def Delete(self):
        self.owner.items.remove(self)


class FakeShapes(FakeCollection):

    def __init__(self, items=None, picture_failure=False):
        super().__init__(items)
        self.picture_failure = picture_failure
        self.picture_calls = []

        for index, shape in enumerate(self.items, start=1):
            shape.owner = self
            shape.ZOrderPosition = index

    def AddShape(self, shape_type, left, top, width, height):
        shape = FakeShape(
            "Rectangle",
            left=left,
            top=top,
            width=width,
            height=height,
        )
        self._append(shape)
        return shape

    def AddPicture(self, **kwargs):
        if self.picture_failure:
            raise RuntimeError("PowerPoint could not add pen PNG")

        self.picture_calls.append(kwargs)
        shape = FakeShape(
            "Picture",
            left=kwargs["Left"],
            top=kwargs["Top"],
            width=360.0,
            height=257.0,
        )
        self._append(shape)
        return shape

    def _append(self, shape):
        shape.owner = self
        shape.ZOrderPosition = len(self.items) + 1
        self.items.append(shape)


class FakeSlide:

    def __init__(
        self,
        shapes,
        sequence,
        slide_width=960.0,
        picture_failure=False,
    ):
        self.Shapes = FakeShapes(
            shapes,
            picture_failure=picture_failure,
        )
        self.TimeLine = SimpleNamespace(
            MainSequence=sequence
        )
        self.SlideShowTransition = SimpleNamespace(
            EntryEffect=0,
            Speed=0,
            AdvanceOnTime=True,
            AdvanceTime=7.5,
        )
        self.Background = SimpleNamespace(
            Fill=FakeFill()
        )
        self.Parent = SimpleNamespace(
            PageSetup=SimpleNamespace(
                SlideWidth=slide_width,
                SlideHeight=540.0,
            )
        )


class FakeSlides:

    def __init__(self, slides):
        self.slides = dict(slides)

    def __call__(self, slide_index):
        return self.slides[slide_index]


class FakeSavedPowerPointController:

    def __init__(self, slides):
        self.presentation = SimpleNamespace(
            Slides=FakeSlides(slides)
        )

    def __enter__(self):
        return self

    def __exit__(self, exc_type, exc_value, traceback):
        return False

    def open_presentation(self, pptx_path):
        return None


class FakeTemplateShape:

    def __init__(self, name):
        self.name = name
        self.text = ""
        self.has_text_frame = False
        self.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE


class FakeTemplateSlide:

    def __init__(self, shape_name):
        self.shapes = [FakeTemplateShape(shape_name)]


def audio_effect(
    duration=3.0,
    delay=0.5,
    name="past_sentence",
    shape_id=5,
):
    return FakeEffect(
        shape=SimpleNamespace(
            Name=name,
            Id=shape_id,
        ),
        effect_type=83,
        trigger_type=3,
        delay=delay,
        duration=duration,
    )


class HandwritingPenAnimationTests(unittest.TestCase):

    TEMPLATE_PATH = "templates/vocabulary_template_v3.pptx"

    def setUp(self):
        self.settings = (
            VisualAnimationSettings.from_project_config()
        )
        self.planner = VisualAnimationPlanner(
            settings=self.settings
        )
        self.processor = VisualAnimationPresentationProcessor(
            planner=self.planner
        )

    def _spec(self, slide_within_word=2, planner=None):
        planner = planner or self.planner
        shape_name = {
            2: "PAST_SENTENCE",
            3: "PRESENT_SENTENCE",
            4: "FUTURE_SENTENCE",
        }[slide_within_word]

        return planner.plan_slide(
            FakeTemplateSlide(shape_name),
            slide_within_word,
        )[0]

    def _process(
        self,
        slide_within_word=2,
        sentence=None,
        settings=None,
        picture_failure=False,
        extra_shapes=None,
        audio=None,
        leading_effects=None,
    ):
        settings = settings or self.settings
        planner = VisualAnimationPlanner(settings=settings)
        processor = VisualAnimationPresentationProcessor(
            planner=planner
        )
        shape_name = {
            2: "PAST_SENTENCE",
            3: "PRESENT_SENTENCE",
            4: "FUTURE_SENTENCE",
        }[slide_within_word]
        sentence = sentence or FakeShape(shape_name)
        audio = audio or audio_effect()
        sequence = FakeSequence(
            list(leading_effects or []) + [audio]
        )
        slide = FakeSlide(
            [sentence] + list(extra_shapes or []),
            sequence,
            picture_failure=picture_failure,
        )
        slide.SlideShowTransition.AdvanceTime = (
            audio.Timing.TriggerDelayTime
            + audio.Timing.Duration
        )
        processor.process_slide(
            slide,
            [self._spec(slide_within_word, planner)],
            slide_within_word,
            slide_within_word,
        )

        return processor, slide, sequence, sentence, audio

    @staticmethod
    def _saved_fade_xml(
        shape_id=1,
        target_shape_id=None,
        include_fade=True,
        include_by_letter=True,
        duplicate=False,
        include_mask=False,
        milliseconds=35,
    ):
        target_shape_id = (
            shape_id
            if target_shape_id is None
            else target_shape_id
        )
        iterate = (
            '<p:iterate type="lt">'
            f'<p:tmAbs val="{milliseconds}"/>'
            '</p:iterate>'
            if include_by_letter
            else ""
        )
        preset_id = "10" if include_fade else "22"
        filter_name = "fade" if include_fade else "wipe(left)"
        effect = (
            f'<p:cTn presetID="{preset_id}" presetClass="entr">'
            f"{iterate}"
            '<p:animEffect transition="in" '
            f'filter="{filter_name}">'
            f'<p:spTgt spid="{target_shape_id}"/>'
            '</p:animEffect>'
            '</p:cTn>'
        )
        mask = (
            '<p:cNvPr id="99" '
            'name="PAST_SENTENCE_REVEAL_MASK"/>'
            if include_mask
            else ""
        )
        return (
            '<p:sld xmlns:p="http://schemas.openxmlformats.org/'
            'presentationml/2006/main">'
            f'<p:cNvPr id="{shape_id}" name="PAST_SENTENCE"/>'
            f"{mask}{effect}"
            f"{effect if duplicate else ''}"
            '</p:sld>'
        )

    def test_handwriting_pen_plan_uses_text_only_fallback(self):
        spec = self._spec(2)

        self.assertEqual(spec.reveal_mode, "handwriting_pen")
        self.assertEqual(spec.effect_id, 0)
        self.assertIsNotNone(spec.handwriting_pen)
        self.assertEqual(
            spec.handwriting_pen.fallback_effect,
            "text_only",
        )
        self.assertEqual(
            spec.handwriting_pen.letter_delay,
            project_config.ANIMATION_HANDWRITING_LETTER_DELAY,
        )
        self.assertEqual(
            spec.handwriting_pen.audio_gap,
            project_config.ANIMATION_HANDWRITING_AUDIO_GAP,
        )

    def test_pen_asset_path_comes_from_config(self):
        spec = self._spec(2)

        self.assertEqual(
            spec.handwriting_pen.image_path,
            Path(
                project_config.ANIMATION_HANDWRITING_PEN_IMAGE
            ).resolve(),
        )

    def test_pen_geometry_is_derived_from_rendered_text_bounds(self):
        sentence = FakeShape(
            "PAST_SENTENCE",
            left=145.0,
            top=310.0,
            width=510.0,
            height=32.0,
            line_bounds=[(180.0, 316.0, 155.0, 20.0)],
        )
        _, slide, _, _, _ = self._process(
            sentence=sentence
        )
        pen = slide.Shapes.items[-1]

        self.assertEqual(
            pen.Name,
            "PAST_SENTENCE_HANDWRITING_PEN",
        )
        self.assertEqual(
            pen.Left,
            180.0
            + self.settings.handwriting_pen_offset_x,
        )
        self.assertEqual(
            pen.Top,
            316.0
            + self.settings.handwriting_pen_offset_y,
        )
        self.assertEqual(
            pen.Width,
            self.settings.handwriting_pen_width,
        )

    def test_handwriting_distance_uses_text_not_textbox_width(self):
        sentence = FakeShape(
            "PAST_SENTENCE",
            left=100.0,
            width=760.0,
            line_bounds=[(125.0, 204.0, 140.0, 20.0)],
        )
        _, slide, sequence, _, _ = self._process(
            sentence=sentence
        )
        pen_effect = sequence.effects[1]
        pen_motion = (
            pen_effect.Behaviors.items[0].MotionEffect
        )
        expected_pen = 140.0 / 960.0 * 100.0

        self.assertAlmostEqual(pen_motion.ByX, expected_pen)
        self.assertEqual(pen_motion.ByY, 0.0)
        pen = slide.Shapes.items[-1]
        pen_end_left = (
            pen.Left
            + pen_motion.ByX / 100.0 * 960.0
        )
        self.assertGreaterEqual(
            pen_end_left,
            125.0
            + self.settings.handwriting_pen_offset_x
            + 140.0,
        )
        self.assertLess(pen_end_left, sentence.Left + sentence.Width)

    def test_fade_by_letter_and_pen_start_together(self):
        audio = audio_effect(duration=3.0, delay=0.7)
        _, _, sequence, sentence, audio = self._process(audio=audio)
        text_effect = sequence.effects[0]
        pen_effect = sequence.effects[1]

        self.assertIs(sequence.effects[-1], audio)
        self.assertEqual(
            text_effect.EffectType,
            self.processor.FADE_EFFECT,
        )
        self.assertEqual(
            text_effect.EffectInformation.TextUnitEffect,
            self.processor.TEXT_UNIT_BY_CHARACTER,
        )
        self.assertTrue(sentence.AnimationSettings.Animate)
        self.assertEqual(
            sentence.AnimationSettings.TextUnitEffect,
            self.processor.PP_ANIMATE_BY_CHARACTER,
        )
        self.assertNotEqual(
            sentence.AnimationSettings.TextLevelEffect,
            self.processor.PP_ANIMATE_LEVEL_NONE,
        )
        self.assertEqual(
            pen_effect.Timing.TriggerType,
            self.processor.WITH_PREVIOUS,
        )
        self.assertEqual(
            text_effect.Timing.TriggerType,
            self.processor.AFTER_PREVIOUS,
        )
        self.assertEqual(
            text_effect.Timing.TriggerDelayTime,
            0.7 + self.settings.sentence_delay,
        )
        self.assertEqual(
            pen_effect.Timing.TriggerDelayTime,
            text_effect.Timing.TriggerDelayTime,
        )
        self.assertEqual(
            text_effect.Timing.Duration,
            self.settings.handwriting_letter_delay,
        )

    def test_slides_2_3_and_4_target_only_semantic_sentence(self):
        plans = self.planner.build_template_plan(
            self.TEMPLATE_PATH
        )

        self.assertEqual(
            [
                [spec.shape_name for spec in plans[index]]
                for index in (1, 2, 3)
            ],
            [
                ["PAST_SENTENCE"],
                ["PRESENT_SENTENCE"],
                ["FUTURE_SENTENCE"],
            ],
        )
        self.assertTrue(
            all(
                spec.reveal_mode == "handwriting_pen"
                for plan in plans[1:]
                for spec in plan
            )
        )

    def test_multiline_sentence_has_left_to_right_write_segments(self):
        sentence = FakeShape(
            "PAST_SENTENCE",
            width=700.0,
            height=70.0,
            line_bounds=[
                (120.0, 200.0, 300.0, 20.0),
                (125.0, 235.0, 180.0, 20.0),
            ],
        )
        processor, slide, sequence, _, _ = self._process(
            sentence=sentence
        )

        self.assertEqual(processor.handwriting_pen_count, 1)
        self.assertEqual(processor.handwriting_fallback_count, 0)
        self.assertEqual(len(sequence.effects), 6)
        self.assertEqual(slide.Shapes.Count, 2)
        motions = [
            effect.Behaviors.items[0].MotionEffect
            for effect in sequence.effects[1:4]
        ]
        self.assertGreater(motions[0].ByX, 0.0)
        self.assertEqual(motions[0].ByY, 0.0)
        self.assertLess(motions[1].ByX, 0.0)
        self.assertGreater(motions[1].ByY, 0.0)
        self.assertGreater(motions[2].ByX, 0.0)
        self.assertEqual(motions[2].ByY, 0.0)

    def test_line_return_moves_to_next_rendered_line_start(self):
        lines = [
            (120.0, 200.0, 300.0, 20.0),
            (125.0, 235.0, 180.0, 20.0),
        ]
        sentence = FakeShape(
            "PAST_SENTENCE",
            line_bounds=lines,
        )
        processor, _, sequence, _, _ = self._process(
            sentence=sentence
        )
        return_motion = (
            sequence.effects[2]
            .Behaviors.items[0]
            .MotionEffect
        )

        self.assertAlmostEqual(
            return_motion.ByX,
            (125.0 - (120.0 + 300.0)) / 960.0 * 100.0,
        )
        self.assertAlmostEqual(
            return_motion.ByY,
            (235.0 - 200.0) / 540.0 * 100.0,
        )

    def test_line_durations_are_proportional_and_bounded(self):
        sentence = FakeShape(
            "PAST_SENTENCE",
            line_bounds=[
                (120.0, 200.0, 150.0, 20.0),
                (125.0, 235.0, 300.0, 20.0),
            ],
            line_texts=[
                "twenty visible letters here",
                "short line",
            ],
        )
        processor, _, sequence, _, _ = self._process(
            sentence=sentence
        )
        write_one = sequence.effects[1].Timing.Duration
        line_return = sequence.effects[2].Timing.Duration
        write_two = sequence.effects[3].Timing.Duration

        first_visible_characters = 27
        second_visible_characters = 10
        self.assertAlmostEqual(
            write_one / write_two,
            first_visible_characters
            / second_visible_characters,
        )
        self.assertGreater(write_one, write_two)
        self.assertLess(
            sentence.TextFrame2.TextRange.Lines.Item(1).BoundWidth,
            sentence.TextFrame2.TextRange.Lines.Item(2).BoundWidth,
        )
        self.assertEqual(
            line_return,
            self.settings.handwriting_line_return_duration,
        )
        self.assertLess(line_return, write_one)
        expected_total = (
            (
                first_visible_characters
                + second_visible_characters
            )
            * self.settings.handwriting_letter_delay
            + self.settings.handwriting_line_return_duration
        )
        self.assertAlmostEqual(
            write_one + line_return + write_two,
            expected_total,
        )
        timing_plan = processor.handwriting_timing_plans[2]
        self.assertEqual(
            timing_plan.lines[0].writing_duration,
            write_one,
        )
        self.assertEqual(
            timing_plan.lines[1].writing_duration,
            write_two,
        )
        self.assertEqual(
            timing_plan.lines[1].start_time,
            timing_plan.lines[0].end_time,
        )

    def test_handwriting_pen_creates_no_reveal_masks(self):
        sentence = FakeShape(
            "PAST_SENTENCE",
            width=700.0,
            line_bounds=[
                (120.0, 200.0, 300.0, 20.0),
                (125.0, 235.0, 180.0, 20.0),
            ],
        )
        _, slide, _, _, _ = self._process(sentence=sentence)
        masks = [
            shape
            for shape in slide.Shapes.items
            if "_REVEAL_MASK_LINE_" in shape.Name
        ]

        self.assertEqual(masks, [])
        self.assertTrue(
            all(
                "REVEAL_MASK" not in shape.Name
                for shape in slide.Shapes.items
            )
        )

    def test_configured_letter_delay_serializes_as_tmabs(self):
        xml = (
            '<p:sld xmlns:p="http://schemas.openxmlformats.org/'
            'presentationml/2006/main">'
            '<p:cNvPr id="7" name="PAST_SENTENCE"/>'
            '<p:cTn presetID="10" presetClass="entr">'
            '<p:iterate type="lt"><p:tmPct val="10000"/>'
            '</p:iterate>'
            '<p:animEffect transition="in" filter="fade">'
            '<p:spTgt spid="7"/>'
            '</p:animEffect>'
            '</p:cTn>'
            '<p:cNvPr id="8" name="UNRELATED"/>'
            '<p:cTn presetID="10" presetClass="entr">'
            '<p:iterate type="lt"><p:tmPct val="20000"/>'
            '</p:iterate>'
            '<p:animEffect transition="in" filter="fade">'
            '<p:spTgt spid="8"/>'
            '</p:animEffect>'
            '</p:cTn>'
            '</p:sld>'
        )
        milliseconds = round(
            self.settings.handwriting_letter_delay * 1000
        )
        patched = SentenceLetterTimingPatcher().patch_slide_xml(
            xml,
            "PAST_SENTENCE",
            milliseconds,
        )

        self.assertIn(
            f'<p:tmAbs val="{milliseconds}"/>',
            patched,
        )
        self.assertIn('<p:tmPct val="20000"/>', patched)
        self.assertEqual(patched.count('<p:tmAbs '), 1)

    def test_saved_xml_preserves_fade_by_letter_target_and_timing(self):
        milliseconds = round(
            self.settings.handwriting_letter_delay * 1000
        )
        state = SentenceLetterTimingPatcher().verify_slide_xml(
            self._saved_fade_xml(
                shape_id=15,
                milliseconds=milliseconds,
            ),
            "PAST_SENTENCE",
            milliseconds,
            expected_shape_id=15,
        )

        self.assertEqual(state.shape_name, "PAST_SENTENCE")
        self.assertEqual(state.shape_id, 15)
        self.assertEqual(state.letter_delay_ms, milliseconds)

    def test_saved_xml_rejects_genuinely_missing_fade(self):
        with self.assertRaises(AnimationTemplateError):
            SentenceLetterTimingPatcher().verify_slide_xml(
                self._saved_fade_xml(include_fade=False),
                "PAST_SENTENCE",
                35,
            )

    def test_saved_xml_rejects_missing_by_letter_iteration(self):
        with self.assertRaises(AnimationTemplateError):
            SentenceLetterTimingPatcher().verify_slide_xml(
                self._saved_fade_xml(include_by_letter=False),
                "PAST_SENTENCE",
                35,
            )

    def test_saved_xml_rejects_wrong_target_shape(self):
        with self.assertRaises(AnimationTemplateError):
            SentenceLetterTimingPatcher().verify_slide_xml(
                self._saved_fade_xml(target_shape_id=2),
                "PAST_SENTENCE",
                35,
            )

    def test_saved_xml_rejects_duplicate_sentence_effects(self):
        with self.assertRaises(AnimationTemplateError):
            SentenceLetterTimingPatcher().verify_slide_xml(
                self._saved_fade_xml(duplicate=True),
                "PAST_SENTENCE",
                35,
            )

    def test_saved_xml_rejects_reveal_masks(self):
        with self.assertRaises(AnimationTemplateError):
            SentenceLetterTimingPatcher().verify_slide_xml(
                self._saved_fade_xml(include_mask=True),
                "PAST_SENTENCE",
                35,
            )

    def test_saved_verifier_accepts_normalized_com_with_video_effect(self):
        processor, slide, sequence, sentence, _audio = (
            self._process()
        )
        video = FakeEffect(
            shape=SimpleNamespace(
                Name="VOCAB_IMAGE",
                Id=6,
            ),
            effect_type=processor.MEDIA_PLAY_EFFECT,
            trigger_type=processor.WITH_PREVIOUS,
            duration=0.001,
            sequence=sequence,
        )
        video.Timing.TriggerDelayTime = 0.0
        sequence.effects[0].Timing.TriggerType = (
            processor.WITH_PREVIOUS
        )
        sequence.effects.insert(0, video)
        milliseconds = round(
            processor.planner.settings.handwriting_letter_delay
            * 1000
        )
        xml_state = SentenceLetterTimingPatcher().verify_slide_xml(
            self._saved_fade_xml(
                shape_id=sentence.Id,
                milliseconds=milliseconds,
            ),
            "PAST_SENTENCE",
            milliseconds,
            expected_shape_id=sentence.Id,
        )
        processor.controller_factory = (
            lambda visible: FakeSavedPowerPointController(
                {2: slide}
            )
        )

        processor._verify_saved_handwriting(
            "saved.pptx",
            {2: xml_state},
        )

        self.assertEqual(
            sentence.AnimationSettings.EntryEffect,
            257,
        )
        self.assertEqual(
            [
                effect
                for effect in sequence.effects
            if effect.Shape is sentence
            and effect.EffectType == processor.FADE_EFFECT
        ],
            [sequence.effects[1]],
        )

    def test_video_before_audio_does_not_erase_or_duplicate_handwriting(self):
        video = FakeEffect(
            shape=SimpleNamespace(
                Name="VOCAB_IMAGE",
                Id=6,
            ),
            effect_type=self.processor.MEDIA_PLAY_EFFECT,
            trigger_type=self.processor.WITH_PREVIOUS,
            duration=0.001,
        )
        processor, _slide, sequence, sentence, audio = (
            self._process(leading_effects=[video])
        )
        sentence_fades = [
            effect
            for effect in sequence.effects
            if effect.Shape is sentence
            and effect.EffectType == processor.FADE_EFFECT
        ]

        self.assertEqual(len(sentence_fades), 1)
        self.assertEqual(
            sentence_fades[0].Timing.TriggerType,
            processor.WITH_PREVIOUS,
        )
        self.assertIn(video, sequence.effects)
        self.assertIs(sequence.effects[-1], audio)
        self.assertEqual(
            processor.saved_handwriting_targets[2].audio_shape_id,
            audio.Shape.Id,
        )

    def test_video_replacement_preserves_vocab_image_semantic_name(self):
        processor = VideoPresentationProcessor()
        picture = SimpleNamespace(
            Name="VOCAB_IMAGE",
            Left=10,
            Top=20,
            Width=300,
            Height=200,
            ZOrderPosition=4,
            Delete=lambda: None,
        )
        media_shape = SimpleNamespace(
            Name="candidate_video",
            Id=20,
            ZOrderPosition=4,
        )
        processor.locator = SimpleNamespace(
            find_picture=lambda slide: picture
        )
        processor.video_embedder = SimpleNamespace(
            embed=lambda **arguments: media_shape
        )
        processor.video_normalizer = SimpleNamespace(
            prepare=lambda path: Path(path)
        )
        slide = SimpleNamespace(
            SlideShowTransition=SimpleNamespace(
                AdvanceTime=7.5
            ),
            TimeLine=SimpleNamespace(
                MainSequence=SimpleNamespace(
                    Count=0,
                    Item=lambda index: None,
                )
            ),
        )

        processor._replace_picture_with_video(
            slide=slide,
            word=SimpleNamespace(default_video="video.mp4"),
            slide_index=45,
            slide_type="VOCABULARY",
        )

        self.assertEqual(media_shape.Name, "VOCAB_IMAGE")

    def test_pen_fades_out_before_audio(self):
        _, _, sequence, _, audio = self._process()
        hide_effect = sequence.effects[-2]

        self.assertIs(sequence.effects[-1], audio)
        self.assertEqual(
            hide_effect.EffectType,
            self.processor.FADE_EFFECT,
        )
        self.assertEqual(hide_effect.Exit, self.processor.MSO_TRUE)
        self.assertEqual(
            hide_effect.Timing.Duration,
            self.settings.handwriting_pen_hide_duration,
        )

    def test_reveal_mask_mode_remains_independently_available(self):
        settings = replace(
            self.settings,
            sentence_effect="reveal_mask",
        )
        _, slide, sequence, _, _ = self._process(
            settings=settings
        )

        self.assertEqual(len(sequence.effects), 2)
        self.assertEqual(
            slide.Shapes.items[-1].Name,
            "PAST_SENTENCE_REVEAL_MASK",
        )

    def test_missing_pen_asset_uses_text_only_fallback(self):
        settings = replace(
            self.settings,
            handwriting_pen_image=(
                "presentation/assets/handwriting/missing.png"
            ),
        )
        output = io.StringIO()

        with redirect_stdout(output):
            processor, slide, sequence, _, audio = self._process(
                settings=settings
            )

        self.assertEqual(processor.handwriting_fallback_count, 1)
        self.assertEqual(len(sequence.effects), 2)
        self.assertEqual(slide.Shapes.Count, 1)
        self.assertIn("missing or invalid", output.getvalue())
        self.assertNotIn("REVEAL_MASK", output.getvalue())
        self.assertIs(sequence.effects[-1], audio)

    def test_corrupt_pen_asset_uses_text_only_fallback(self):
        with tempfile.TemporaryDirectory() as temp_folder:
            corrupt_path = Path(temp_folder) / "hand_pen.png"
            corrupt_path.write_bytes(b"not a PNG")
            settings = replace(
                self.settings,
                handwriting_pen_image=str(corrupt_path),
            )
            output = io.StringIO()

            with redirect_stdout(output):
                processor, slide, sequence, _, _ = self._process(
                    settings=settings
                )

        self.assertEqual(processor.handwriting_fallback_count, 1)
        self.assertEqual(len(sequence.effects), 2)
        self.assertEqual(slide.Shapes.Count, 1)
        self.assertIn("missing or invalid", output.getvalue())

    def test_picture_failure_uses_text_only_fallback(self):
        output = io.StringIO()

        with redirect_stdout(output):
            processor, slide, sequence, _, _ = self._process(
                picture_failure=True
            )

        self.assertEqual(processor.handwriting_fallback_count, 1)
        self.assertEqual(len(sequence.effects), 2)
        self.assertEqual(slide.Shapes.Count, 1)
        self.assertIn("could not add pen PNG", output.getvalue())

    def test_audio_starts_after_handwriting_gap_and_keeps_duration(self):
        audio = audio_effect(duration=2.7, delay=0.5)
        original_duration = audio.Timing.Duration
        _, slide, sequence, sentence, audio = self._process(audio=audio)
        visible_character_count = len(
            sentence.TextFrame2.TextRange.Text
            .replace("\r", "")
            .replace("\n", "")
            .replace("\v", "")
        )
        handwriting_duration = (
            visible_character_count
            * self.settings.handwriting_letter_delay
            + self.settings.handwriting_pen_hide_duration
        )
        expected_start = (
            0.5
            + self.settings.sentence_delay
            + handwriting_duration
            + self.settings.handwriting_audio_gap
        )

        self.assertIs(sequence.effects[-1], audio)
        self.assertEqual(
            audio.Timing.TriggerDelayTime,
            self.settings.handwriting_audio_gap,
        )
        self.assertEqual(audio.Timing.Duration, original_duration)
        self.assertAlmostEqual(
            slide.SlideShowTransition.AdvanceTime,
            expected_start + original_duration,
        )

    def test_slide_end_preserves_existing_end_padding(self):
        audio = audio_effect(duration=2.0, delay=0.5)
        sentence = FakeShape("PAST_SENTENCE")
        sequence = FakeSequence([audio])
        slide = FakeSlide([sentence], sequence)
        slide.SlideShowTransition.AdvanceTime = 2.8
        planner = VisualAnimationPlanner(settings=self.settings)
        processor = VisualAnimationPresentationProcessor(planner=planner)
        processor.process_slide(
            slide,
            [self._spec(2, planner)],
            2,
            2,
        )
        new_audio_start = (
            0.5
            + len(
                sentence.TextFrame2.TextRange.Text
                .replace("\r", "")
                .replace("\n", "")
                .replace("\v", "")
            )
            * self.settings.handwriting_letter_delay
            + self.settings.handwriting_pen_hide_duration
            + self.settings.handwriting_audio_gap
        )

        self.assertEqual(
            slide.SlideShowTransition.AdvanceTime,
            new_audio_start + 2.0 + 0.3,
        )

    def test_pen_normalization_preserves_rgba_and_source(self):
        with tempfile.TemporaryDirectory() as temp_folder:
            source_path = Path(temp_folder) / "hand_pen.png"
            image = Image.new("RGBA", (7, 7))
            pixels = image.load()

            for y in range(7):
                for x in range(7):
                    background = (
                        (255, 255, 255, 255)
                        if (x + y) % 2 == 0
                        else (239, 239, 238, 255)
                    )
                    pixels[x, y] = background

            for y in range(2, 5):
                for x in range(2, 5):
                    pixels[x, y] = (190, 110, 80, 255)

            pixels[3, 3] = (255, 255, 255, 255)
            image.save(source_path, format="PNG")
            source_bytes = source_path.read_bytes()
            normalizer = HandwritingPenAssetNormalizer()

            normalized_path = normalizer.prepare(
                source_path,
                alpha_threshold=8,
                background_tolerance=2,
            )
            cached_path = normalizer.prepare(
                source_path,
                alpha_threshold=8,
                background_tolerance=2,
            )

            self.assertNotEqual(normalized_path, source_path)
            self.assertEqual(normalized_path, cached_path)
            self.assertEqual(source_path.read_bytes(), source_bytes)

            with Image.open(normalized_path) as normalized:
                normalized.load()
                self.assertEqual(normalized.mode, "RGBA")
                self.assertEqual(
                    normalized.getchannel("A").getextrema(),
                    (0, 255),
                )

    def test_progress_shapes_remain_untouched(self):
        track = FakeShape(
            "PROGRESS_TRACK",
            left=20,
            top=520,
            width=900,
            height=8,
        )
        fill = FakeShape(
            "PROGRESS_FILL",
            left=20,
            top=520,
            width=450,
            height=8,
        )
        before = [
            (shape.Left, shape.Top, shape.Width, shape.Height)
            for shape in (track, fill)
        ]
        _, _, sequence, _, _ = self._process(
            extra_shapes=[track, fill]
        )

        self.assertEqual(
            [
                (shape.Left, shape.Top, shape.Width, shape.Height)
                for shape in (track, fill)
            ],
            before,
        )
        self.assertTrue(
            all(
                effect.Shape not in (track, fill)
                for effect in sequence.effects
                if effect.Shape is not None
            )
        )

    def test_normal_logging_has_one_handwriting_summary(self):
        self.processor.handwriting_pen_count = 3
        self.processor.handwriting_fallback_count = 1
        output = io.StringIO()

        with redirect_stdout(output):
            self.processor._log_handwriting_summary()

        self.assertEqual(
            output.getvalue().strip(),
            "Handwriting pen... 3 sentences, "
            "1 text-only fallbacks",
        )
        self.assertNotIn("motion=", output.getvalue())

    def test_verbose_logging_includes_pen_motion_detail(self):
        output = io.StringIO()

        with patch.object(
            project_config,
            "PRESENTATION_VERBOSE_LOGGING",
            True,
        ), redirect_stdout(output):
            self._process()

        self.assertIn("handwriting pen", output.getvalue())
        self.assertIn("audio=", output.getvalue())


if __name__ == "__main__":
    unittest.main()
