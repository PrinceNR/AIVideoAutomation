

from pathlib import Path
from pptx import Presentation

prs = Presentation("templates/vocabulary_template_v2.pptx")

for i, slide in enumerate(prs.slides, start=1):

    print(f"\nSlide {i}")

    for shape in slide.shapes:

        print(
            shape.name,
            shape.shape_type,
            shape.left,
            shape.top,
            shape.width,
            shape.height
        )