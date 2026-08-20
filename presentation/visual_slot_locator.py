from pptx.enum.shapes import MSO_SHAPE_TYPE


class VisualSlotLocator:

    IMAGE_SHAPE_NAME = "VOCAB_IMAGE"

    # Legacy template fallback
    IMAGE_LEFT = 1223367
    IMAGE_TOP = 728662
    IMAGE_WIDTH = 3344465

    def find_picture(
        self,
        slide
    ):

        # ---------------------------------
        # Preferred method: semantic name
        # ---------------------------------

        for shape in slide.shapes:

            if (
                shape.shape_type
                == MSO_SHAPE_TYPE.PICTURE
                and shape.name
                == self.IMAGE_SHAPE_NAME
            ):
                return shape

        # ---------------------------------
        # Legacy fallback: old geometry
        # ---------------------------------

        for shape in slide.shapes:

            if (
                shape.shape_type
                == MSO_SHAPE_TYPE.PICTURE
                and shape.left
                == self.IMAGE_LEFT
                and shape.top
                == self.IMAGE_TOP
                and shape.width
                == self.IMAGE_WIDTH
            ):
                return shape

        return None