from models.word import Word
from pptx.slide import Slide
from pptx.enum.shapes import MSO_SHAPE_TYPE


class TextProcessor:

    def process(
        self,
        slide,
        slide_definition,
        word,
        word_number,
        total_words,
        timeline
    ):

        placeholders = self._get_word_placeholders(
            word,
            word_number,
            total_words
        )

        for placeholder, value in placeholders.items():

            self.replace_text(
                slide,
                placeholder,
                str(value)
            )


    def _get_word_placeholders(
        self,
        word: Word,
        word_number: int,
        total_words: int
    ) -> dict:

        return {
            "{{WORD}}": word.word,
            "{{MEANING}}": word.meaning,
            "{{PRONUNCIATION}}": word.pronunciation,

            "{{MALAYALAM}}": word.translations.get(
                "malayalam",
                ""
            ),

            "{{TAMIL}}": word.translations.get(
                "tamil",
                ""
            ),

            "{{HINDI}}": word.translations.get(
                "hindi",
                ""
            ),

            "{{PRESENT_SENTENCE}}":
                word.present_sentence,

            "{{PAST_SENTENCE}}":
                word.past_sentence,

            "{{FUTURE_SENTENCE}}":
                word.future_sentence,

            "{{BASE_FORM}}":
                word.base_form,

            "{{PRESENT_FORM}}":
                word.present_form,

            "{{PAST_FORM}}":
                word.past_form,

            "{{PART_OF_SPEECH}}":
                word.part_of_speech,

            "{{DIFFICULTY}}":
                word.difficulty,

            "{{WORD_NUMBER}}":
                f"{word_number}/{total_words}",
        }


    def replace_text(
        self,
        slide: Slide,
        placeholder: str,
        value: str
    ):

        for shape in self._iter_shapes(
            slide.shapes
        ):

            if not shape.has_text_frame:
                continue

            for paragraph in (
                shape.text_frame.paragraphs
            ):

                self._replace_in_paragraph(
                    paragraph,
                    placeholder,
                    value
                )


    def _iter_shapes(
        self,
        shapes
    ):

        for shape in shapes:

            # Also support text boxes
            # inside grouped PowerPoint objects.
            if (
                shape.shape_type
                == MSO_SHAPE_TYPE.GROUP
            ):

                yield from self._iter_shapes(
                    shape.shapes
                )

            else:

                yield shape


    def _replace_in_paragraph(
        self,
        paragraph,
        placeholder: str,
        value: str
    ):

        runs = list(
            paragraph.runs
        )

        if not runs:
            return

        # ---------------------------------
        # Case 1:
        # placeholder exists inside
        # a normal single run
        # ---------------------------------

        for run in runs:

            if placeholder in run.text:

                run.text = run.text.replace(
                    placeholder,
                    value
                )

        # ---------------------------------
        # Case 2:
        # PowerPoint split placeholder
        # across multiple runs
        # ---------------------------------

        while True:

            runs = list(
                paragraph.runs
            )

            full_text = "".join(
                run.text
                for run in runs
            )

            start = full_text.find(
                placeholder
            )

            if start == -1:
                break

            end = (
                start
                + len(placeholder)
            )

            start_run_index = None
            end_run_index = None

            start_offset = None
            end_offset = None

            cursor = 0

            for index, run in enumerate(
                runs
            ):

                run_start = cursor
                run_end = (
                    cursor
                    + len(run.text)
                )

                if (
                    start_run_index is None
                    and start < run_end
                ):

                    start_run_index = index

                    start_offset = (
                        start
                        - run_start
                    )

                if (
                    end_run_index is None
                    and end <= run_end
                    and end > run_start
                ):

                    end_run_index = index

                    end_offset = (
                        end
                        - run_start
                    )

                    break

                cursor = run_end

            if (
                start_run_index is None
                or end_run_index is None
            ):
                break

            start_run = runs[
                start_run_index
            ]

            end_run = runs[
                end_run_index
            ]

            prefix = (
                start_run.text[
                    :start_offset
                ]
            )

            suffix = (
                end_run.text[
                    end_offset:
                ]
            )

            if (
                start_run_index
                == end_run_index
            ):

                start_run.text = (
                    prefix
                    + value
                    + suffix
                )

            else:

                # Put replacement into
                # the first run so it keeps
                # that run's formatting.
                start_run.text = (
                    prefix
                    + value
                )

                # Remove fragments of the
                # placeholder from middle runs.
                for index in range(
                    start_run_index + 1,
                    end_run_index
                ):

                    runs[index].text = ""

                # Keep any text that existed
                # after the placeholder.
                end_run.text = suffix