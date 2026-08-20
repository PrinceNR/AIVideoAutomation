from copy import deepcopy
from pptx.enum.shapes import MSO_SHAPE_TYPE
from presentation.group_copier import GroupCopier


class SlideDuplicator:


    def __init__(self):

        self.group_copier = GroupCopier()

    def duplicate_slide(
        self,
        presentation,
        slide_index
    ):

        source_slide = presentation.slides[slide_index]

        new_slide = presentation.slides.add_slide(
            source_slide.slide_layout
        )

        self.group_copier.copy_shapes(
            source_slide,
            new_slide
        )

        return new_slide