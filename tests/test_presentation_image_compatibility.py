import tempfile
import unittest
from pathlib import Path

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

from presentation.image_compatibility import (
    PresentationImageCompatibility,
    PresentationImageError,
)


class PresentationImageCompatibilityTests(unittest.TestCase):
    def setUp(self):
        self.temporary_directory = tempfile.TemporaryDirectory()
        self.directory = Path(self.temporary_directory.name)
        self.compatibility = PresentationImageCompatibility()

    def tearDown(self):
        self.temporary_directory.cleanup()

    def _create_image(self, name, image_format):
        path = self.directory / name
        Image.new("RGBA", (20, 20), (255, 0, 0, 128)).save(
            path,
            format=image_format,
        )
        return path

    def test_supported_jpeg_and_png_pass_through_unchanged(self):
        jpeg = self.directory / "photo.jpg"
        Image.new("RGB", (20, 20), "red").save(jpeg, format="JPEG")
        png = self._create_image("illustration.png", "PNG")

        self.assertEqual(jpeg, self.compatibility.prepare(jpeg))
        self.assertEqual(png, self.compatibility.prepare(png))

    def test_webp_converts_to_png_usable_by_python_pptx(self):
        webp = self._create_image("manual.webp", "WEBP")

        converted = self.compatibility.prepare(webp)

        self.assertEqual(".png", converted.suffix)
        self.assertTrue(converted.is_file())
        presentation = Presentation()
        slide = presentation.slides.add_slide(
            presentation.slide_layouts[6]
        )
        picture = slide.shapes.add_picture(
            str(converted), 0, 0, Inches(1), Inches(1)
        )
        self.assertIsNotNone(picture)
        self.assertTrue(webp.is_file())

    def test_converted_webp_is_reused_on_rerun(self):
        webp = self._create_image("manual.webp", "WEBP")
        first = self.compatibility.prepare(webp)
        first_timestamp = first.stat().st_mtime_ns

        second = self.compatibility.prepare(webp)

        self.assertEqual(first, second)
        self.assertEqual(first_timestamp, second.stat().st_mtime_ns)

    def test_corrupt_webp_raises_clear_controlled_error(self):
        webp = self.directory / "broken.webp"
        webp.write_bytes(b"not a webp image")

        with self.assertRaises(PresentationImageError) as context:
            self.compatibility.prepare(webp)

        self.assertIn(str(webp), str(context.exception))
        self.assertIn("Could not prepare presentation image", str(context.exception))


if __name__ == "__main__":
    unittest.main()
