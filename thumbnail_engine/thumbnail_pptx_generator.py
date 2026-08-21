from pathlib import Path
import re
import tempfile

from PIL import Image
from pptx import Presentation
from config import THUMBNAIL_MAX_WORDS
from presentation.image_compatibility import (
    PresentationImageCompatibility
)


class ThumbnailPptxGenerator:

    def __init__(self):

        self.image_compatibility = (
            PresentationImageCompatibility()
        )

    def generate(
        self,
        lesson,
        lesson_folder,
        template_path,
        output_path
    ):

        lesson_folder = Path(lesson_folder)
        template_path = Path(template_path)
        output_path = Path(output_path)

        if not template_path.exists():
            raise FileNotFoundError(
                f"Thumbnail template not found: {template_path}"
            )

        presentation = Presentation(template_path)

        if len(presentation.slides) == 0:
            raise ValueError(
                "Thumbnail template does not contain a slide."
            )

        slide = presentation.slides[0]

        print(
            f"Words available: {len(lesson.words)}"
        )

        words = lesson.words[:THUMBNAIL_MAX_WORDS]

        with tempfile.TemporaryDirectory() as temp_dir:

            temp_dir = Path(temp_dir)

            for index in range(1, THUMBNAIL_MAX_WORDS + 1):

                image_placeholder = self._find_shape(
                    slide,
                    f"IMAGE_{index}"
                )

                word_placeholder = self._find_shape(
                    slide,
                    f"WORD_{index}"
                )

                if image_placeholder is None:
                    raise ValueError(
                        f"IMAGE_{index} not found in "
                        f"thumbnail template."
                    )

                if word_placeholder is None:
                    raise ValueError(
                        f"WORD_{index} not found in "
                        f"thumbnail template."
                    )

                # -----------------------------------------
                # Slot has a vocabulary word
                # -----------------------------------------

                if index <= len(words):

                    word = words[index - 1]

                    print(
                        f"Slot {index}: {word.word}"
                    )

                    # Replace sample text
                    self._replace_word_text(
                        word_placeholder,
                        word.word
                    )

                    # Use the image selected for this lesson word.
                    image_path = self._find_word_image(
                        word
                    )

                    print(
                        f"  Image: {image_path}"
                    )

                    (
                        prepared_image,
                        image_left,
                        image_top,
                        image_width,
                        image_height
                    ) = self._prepare_image(
                        image_path=image_path,
                        placeholder=image_placeholder
                    )

                    # Center the complete image inside IMAGE_X.
                    slide.shapes.add_picture(
                        str(prepared_image),
                        image_left,
                        image_top,
                        image_width,
                        image_height
                    )

                # -----------------------------------------
                # Unused slot
                # -----------------------------------------

                else:

                    # Remove sample word from template
                    self._replace_word_text(
                        word_placeholder,
                        ""
                    )

        # ---------------------------------------------
        # Save editable PPTX
        # ---------------------------------------------

        output_path.parent.mkdir(
            parents=True,
            exist_ok=True
        )

        presentation.save(output_path)

        print(
            f"Editable thumbnail PPTX created: "
            f"{output_path}"
        )

        return output_path

    # =================================================
    # Find shape by PowerPoint Selection Pane name
    # =================================================

    def _find_shape(self, slide, shape_name):

        for shape in slide.shapes:

            if shape.name == shape_name:
                return shape

        return None

    # =================================================
    # Replace word while trying to preserve template
    # formatting
    # =================================================

    def _replace_word_text(
        self,
        shape,
        text
    ):

        if not shape.has_text_frame:
            raise ValueError(
                f"{shape.name} is not a text box."
            )

        text_frame = shape.text_frame

        # Preserve the existing first run's formatting
        # whenever possible.
        if (
            text_frame.paragraphs
            and text_frame.paragraphs[0].runs
        ):

            first_paragraph = (
                text_frame.paragraphs[0]
            )

            first_run = (
                first_paragraph.runs[0]
            )

            first_run.text = text

            # Clear any additional runs
            for run in first_paragraph.runs[1:]:
                run.text = ""

            # Clear additional paragraphs
            for paragraph in text_frame.paragraphs[1:]:

                for run in paragraph.runs:
                    run.text = ""

        else:

            shape.text = text

    # =================================================
    # Find first image belonging to word
    # =================================================

    def _find_word_image(
        self,
        word
    ):
        if not word.default_image:
            raise FileNotFoundError(
                f"No default image is selected for thumbnail word "
                f"'{word.word}'."
            )

        image_path = Path(word.default_image)

        if not image_path.is_file():
            raise FileNotFoundError(
                f"Selected thumbnail image not found for word "
                f"'{word.word}': {image_path}"
            )

        return image_path

    # =================================================
    # Fit image inside placeholder without cropping
    # =================================================

    def _prepare_image(
        self,
        image_path,
        placeholder
    ):

        prepared_image = (
            self.image_compatibility.prepare(
                image_path
            )
        )

        with Image.open(prepared_image) as image:
            source_width, source_height = image.size

        scale = min(
            placeholder.width / source_width,
            placeholder.height / source_height
        )

        image_width = max(
            1,
            round(source_width * scale)
        )
        image_height = max(
            1,
            round(source_height * scale)
        )

        image_left = (
            placeholder.left
            + (placeholder.width - image_width) // 2
        )
        image_top = (
            placeholder.top
            + (placeholder.height - image_height) // 2
        )

        return (
            prepared_image,
            image_left,
            image_top,
            image_width,
            image_height
        )

    # =================================================
    # Folder-safe word name
    # =================================================

    def _normalize_name(
        self,
        name
    ):

        name = name.strip().lower()

        name = re.sub(
            r"[^a-z0-9]+",
            "_",
            name
        )

        return name.strip("_")
