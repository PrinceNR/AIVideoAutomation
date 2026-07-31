from pptx import Presentation
from models.word import Word


class PlaceholderReplacer:

    def replace_word(
        self,
        presentation: Presentation,
        word: Word,
        word_number: int,
        total_words: int
    ):

        placeholders = self._get_word_placeholders(
            word,
            word_number,
            total_words
        )

        for placeholder, value in placeholders.items():

            self.replace_text(
                presentation,
                placeholder,
                str(value)
            )

    def _get_word_placeholders(
        self,
        word: Word,
        word_number: int,
        total_words: int
    ) -> dict:

        return {
            "{{WORD}}": word.word,
            "{{MEANING}}": word.meaning,
            "{{ABC}}": word.pronunciation,

            "{{MALAYALAM}}": word.translations.get("malayalam", ""),
            "{{TAMIL}}": word.translations.get("tamil", ""),
            "{{HINDI}}": word.translations.get("hindi", ""),

            "{{PRESENT_SENTENCE}}": word.present_sentence,
            "{{PAST_SENTENCE}}": word.past_sentence,
            "{{FUTURE_SENTENCE}}": word.future_sentence,

            "{{BASE_FORM}}": word.base_form,
            "{{PRESENT_FORM}}": word.present_form,
            "{{PAST_FORM}}": word.past_form,

            "{{PART_OF_SPEECH}}": word.part_of_speech,
            "{{DIFFICULTY}}": word.difficulty,

            "{{WORD_NUMBER}}": f"{word_number}/{total_words}",
        }
    def replace_text(
        self,
        presentation: Presentation,
        placeholder: str,
        value: str
    ):
        for slide in presentation.slides:

            for shape in slide.shapes:

                if not shape.has_text_frame:
                    continue

                for paragraph in shape.text_frame.paragraphs:

                    for run in paragraph.runs:

                        if placeholder in run.text:

                            print(f"Replacing {placeholder} -> {value}")

                            run.text = run.text.replace(
                                placeholder,
                                value
                            )