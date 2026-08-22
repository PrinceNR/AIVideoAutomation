import io
import unittest
from contextlib import redirect_stdout
from dataclasses import replace
from unittest.mock import patch

from pptx.enum.shapes import MSO_SHAPE_TYPE

import config as project_config

from presentation.animations.visual_animation_planner import (
    AnimationTemplateError,
    VisualAnimationPlanner,
    VisualAnimationSpec,
)
from presentation.animations.visual_animation_presentation_processor import (
    VisualAnimationPresentationProcessor,
)
from presentation.animations.visual_animation_settings import (
    AnimationConfigurationError,
    VisualAnimationSettings,
)


class FakeTiming:

    def __init__(self, trigger_type, delay, duration):
        self.TriggerType = trigger_type
        self.TriggerDelayTime = delay
        self.Duration = duration


class FakeEffectParameters:

    def __init__(self, effect):
        self.effect = effect
        self._direction = None

    @property
    def Direction(self):
        return self._direction

    @Direction.setter
    def Direction(self, value):
        self._direction = value

        if value is not None:
            self.effect.Timing.Duration = 0.5


class FakeMotionEffect:

    def __init__(self):
        self.ByX = None
        self.ByY = None


class FakeBehavior:

    def __init__(self, behavior_type):
        self.Type = behavior_type
        self.MotionEffect = FakeMotionEffect()


class FakeBehaviors:

    def __init__(self, add_failure=False):
        self.items = []
        self.add_failure = add_failure

    @property
    def Count(self):
        return len(self.items)

    def Add(self, behavior_type):
        if self.add_failure:
            raise RuntimeError("COM motion behavior failed")

        behavior = FakeBehavior(behavior_type)
        self.items.append(behavior)
        return behavior


class FakeEffect:

    def __init__(
        self,
        shape=None,
        effect_id=None,
        trigger_type=3,
        delay=0.0,
        duration=0.0,
        sequence=None,
        behavior_failure=False,
    ):
        self.Shape = shape
        self.EffectType = effect_id
        self.Timing = FakeTiming(
            trigger_type,
            delay,
            duration,
        )
        self.EffectParameters = FakeEffectParameters(self)
        self.Behaviors = FakeBehaviors(behavior_failure)
        self.sequence = sequence

    def Delete(self):
        self.sequence.effects.remove(self)


class FakeSequence:

    def __init__(
        self,
        effects=None,
        add_failure=False,
        behavior_failure=False,
    ):
        self.effects = list(effects or [])
        self.add_failure = add_failure
        self.behavior_failure = behavior_failure

        for effect in self.effects:
            effect.sequence = self

    @property
    def Count(self):
        return len(self.effects)

    def Item(self, index):
        return self.effects[index - 1]

    def AddEffect(
        self,
        shape,
        effect_id,
        level,
        trigger,
        index=None,
    ):
        if self.add_failure:
            raise RuntimeError("COM AddEffect failed")

        effect = FakeEffect(
            shape=shape,
            effect_id=effect_id,
            trigger_type=trigger,
            sequence=self,
            behavior_failure=self.behavior_failure,
        )

        if index is None:
            self.effects.append(effect)
        else:
            self.effects.insert(index - 1, effect)

        return effect


class FakeCollection:

    def __init__(self, items=None):
        self.items = list(items or [])

    @property
    def Count(self):
        return len(self.items)

    def Item(self, index):
        return self.items[index - 1]


class FakeColor:

    def __init__(self, rgb=0xFFFFFF):
        self.RGB = rgb


class FakeFill:

    def __init__(self, fill_type=1, rgb=0xFFFFFF):
        self.Type = fill_type
        self.ForeColor = FakeColor(rgb)
        self.Transparency = 0.0
        self.solid_called = False

    def Solid(self):
        self.solid_called = True
        self.Type = 1


class FakeLine:

    def __init__(self):
        self.Visible = -1


class FakeShape:

    def __init__(
        self,
        name,
        shape_type=1,
        children=None,
        text="",
        left=100.0,
        top=200.0,
        width=320.0,
        height=60.0,
    ):
        self.Name = name
        self.Type = shape_type
        self.GroupItems = FakeCollection(children)
        self.TextFrame = type(
            "FakeTextFrame",
            (),
            {
                "TextRange": type(
                    "FakeTextRange",
                    (),
                    {"Text": text},
                )()
            },
        )()
        self.Left = left
        self.Top = top
        self.Width = width
        self.Height = height
        self.Fill = FakeFill()
        self.Line = FakeLine()
        self.ZOrderPosition = 0
        self.z_order_calls = []
        self.owner = None

    def ZOrder(self, command):
        self.z_order_calls.append(command)
        self.ZOrderPosition = max(
            (
                shape.ZOrderPosition
                for shape in self.owner.items
            ),
            default=0,
        ) + 1

    def Delete(self):
        self.owner.items.remove(self)


class FakeShapes(FakeCollection):

    def __init__(self, items=None, add_failure=False):
        super().__init__(items)
        self.add_failure = add_failure

        for index, shape in enumerate(self.items, start=1):
            shape.owner = self
            shape.ZOrderPosition = index

    def AddShape(self, shape_type, left, top, width, height):
        if self.add_failure:
            raise RuntimeError("COM AddShape failed")

        shape = FakeShape(
            "Rectangle",
            shape_type=shape_type,
            left=left,
            top=top,
            width=width,
            height=height,
        )
        shape.owner = self
        shape.ZOrderPosition = len(self.items) + 1
        self.items.append(shape)
        return shape


class FakeTransition:

    def __init__(self):
        self.EntryEffect = 999
        self.Speed = 2
        self.AdvanceOnTime = True
        self.AdvanceTime = 7.25


class FakeSlide:

    def __init__(
        self,
        shapes=None,
        sequence=None,
        slide_width=960.0,
        background_rgb=0xFDFDFD,
        background_fill_type=1,
        shape_add_failure=False,
    ):
        self.Shapes = FakeShapes(
            shapes,
            add_failure=shape_add_failure,
        )
        self.TimeLine = type(
            "FakeTimeLine",
            (),
            {"MainSequence": sequence or FakeSequence()},
        )()
        self.SlideShowTransition = FakeTransition()
        self.Background = type(
            "FakeBackground",
            (),
            {
                "Fill": FakeFill(
                    background_fill_type,
                    background_rgb,
                )
            },
        )()
        self.Parent = type(
            "FakePresentation",
            (),
            {
                "PageSetup": type(
                    "FakePageSetup",
                    (),
                    {"SlideWidth": slide_width},
                )()
            },
        )()


class FakeTemplateShape:

    def __init__(self, name, text=""):
        self.name = name
        self.text = text
        self.has_text_frame = bool(text)
        self.shape_type = MSO_SHAPE_TYPE.AUTO_SHAPE


class FakeTemplateSlide:

    def __init__(self, shapes=None):
        self.shapes = list(shapes or [])


def audio_effect(duration=2.0, delay=0.5):
    return FakeEffect(
        effect_id=83,
        trigger_type=3,
        delay=delay,
        duration=duration,
    )


class PresentationVisualAnimationTests(unittest.TestCase):

    TEMPLATE_PATH = "templates/vocabulary_template_v3.pptx"

    def setUp(self):
        self.planner = VisualAnimationPlanner()
        self.settings = self.planner.settings
        self.processor = VisualAnimationPresentationProcessor(
            planner=self.planner
        )
        self.template_plans = self.planner.build_template_plan(
            self.TEMPLATE_PATH
        )

    def _sentence_spec(self, slide_within_word=2):
        return self.template_plans[slide_within_word - 1][0]

    def _process_sentence(
        self,
        slide_within_word=2,
        sentence_shape=None,
        sequence=None,
        processor=None,
        spec=None,
        extra_shapes=None,
        **slide_kwargs,
    ):
        semantic_name = {
            2: "PAST_SENTENCE",
            3: "PRESENT_SENTENCE",
            4: "FUTURE_SENTENCE",
        }[slide_within_word]
        sentence_shape = sentence_shape or FakeShape(
            semantic_name,
            text="She paused before entering.",
        )
        sequence = sequence or FakeSequence(
            [audio_effect(duration=3.0)]
        )
        slide = FakeSlide(
            shapes=[sentence_shape] + list(extra_shapes or []),
            sequence=sequence,
            **slide_kwargs,
        )
        processor = processor or self.processor
        processor.process_slide(
            slide,
            [spec or self._sentence_spec(slide_within_word)],
            slide_within_word,
            slide_within_word,
        )
        return slide, sequence, sentence_shape

    def test_slide_1_receives_introductory_animations(self):
        specs = self.template_plans[0]

        self.assertEqual(
            {spec.semantic_element for spec in specs},
            {
                "word",
                "pronunciation",
                "meaning",
                "translation",
                "verb_form",
                "image",
            },
        )
        self.assertTrue(all(not spec.required for spec in specs))

    def test_slide_2_receives_only_past_sentence(self):
        self._assert_sentence_slide(2, "PAST_SENTENCE")

    def test_slide_3_receives_only_present_sentence(self):
        self._assert_sentence_slide(3, "PRESENT_SENTENCE")

    def test_slide_4_receives_only_future_sentence(self):
        self._assert_sentence_slide(4, "FUTURE_SENTENCE")

    def _assert_sentence_slide(self, slide_number, shape_name):
        specs = self.template_plans[slide_number - 1]

        self.assertEqual(len(specs), 1)
        self.assertEqual(specs[0].shape_name, shape_name)
        self.assertEqual(specs[0].reveal_mode, "reveal_mask")
        self.assertEqual(
            specs[0].reveal_direction,
            "left_to_right",
        )
        self.assertEqual(
            specs[0].mask_color,
            self.settings.sentence_mask_rgb(),
        )
        self.assertTrue(specs[0].required)

    def test_animation_policy_reads_reveal_settings_from_config(self):
        with patch.multiple(
            project_config,
            ANIMATION_WORD_EFFECT="fade",
            ANIMATION_WORD_DURATION=0.33,
            ANIMATION_SENTENCE_EFFECT="reveal_mask",
            ANIMATION_SENTENCE_DURATION=0.77,
            ANIMATION_SENTENCE_DIRECTION="left_to_right",
            ANIMATION_SENTENCE_MASK_COLOR="#123456",
        ):
            settings = VisualAnimationSettings.from_project_config()
            planner = VisualAnimationPlanner(settings=settings)

        word_spec = planner.plan_intro_shape(
            "Intro Word",
            "{{WORD}}",
        )
        sentence_spec = planner.plan_slide(
            FakeTemplateSlide(
                [FakeTemplateShape("PAST_SENTENCE")]
            ),
            2,
        )[0]

        self.assertEqual(word_spec.duration, 0.33)
        self.assertEqual(sentence_spec.effect_id, 0)
        self.assertEqual(sentence_spec.duration, 0.77)
        self.assertEqual(sentence_spec.reveal_mode, "reveal_mask")
        self.assertEqual(
            sentence_spec.mask_color,
            settings.sentence_mask_rgb(),
        )

    def test_mask_geometry_matches_sentence_bounds(self):
        sentence = FakeShape(
            "PAST_SENTENCE",
            left=137.0,
            top=248.0,
            width=411.0,
            height=73.0,
        )
        slide, _, _ = self._process_sentence(
            sentence_shape=sentence
        )
        mask = slide.Shapes.items[-1]

        self.assertEqual(
            (mask.Left, mask.Top, mask.Width, mask.Height),
            (sentence.Left, sentence.Top, sentence.Width, sentence.Height),
        )
        self.assertEqual(mask.Name, "PAST_SENTENCE_REVEAL_MASK")

    def test_mask_is_above_sentence_and_uses_slide_background(self):
        slide, _, sentence = self._process_sentence(
            background_rgb=0xA1B2C3
        )
        mask = slide.Shapes.items[-1]

        self.assertGreater(mask.ZOrderPosition, sentence.ZOrderPosition)
        self.assertEqual(
            mask.z_order_calls,
            [self.processor.BRING_TO_FRONT],
        )
        self.assertTrue(mask.Fill.solid_called)
        self.assertEqual(mask.Fill.ForeColor.RGB, 0xA1B2C3)
        self.assertEqual(mask.Line.Visible, self.processor.MSO_FALSE)

    def test_configured_mask_color_is_used_when_background_is_not_solid(self):
        slide, _, _ = self._process_sentence(
            background_fill_type=2
        )
        mask = slide.Shapes.items[-1]

        self.assertEqual(
            mask.Fill.ForeColor.RGB,
            self.settings.sentence_mask_rgb(),
        )

    def test_reveal_distance_is_derived_from_mask_width(self):
        sentence = FakeShape(
            "PAST_SENTENCE",
            width=384.0,
        )
        slide, sequence, _ = self._process_sentence(
            sentence_shape=sentence,
            slide_width=960.0,
        )
        effect = sequence.effects[1]
        motion = effect.Behaviors.items[0].MotionEffect

        self.assertIs(effect.Shape, slide.Shapes.items[-1])
        self.assertEqual(effect.EffectType, self.processor.CUSTOM_EFFECT)
        self.assertEqual(effect.Behaviors.items[0].Type, 1)
        self.assertAlmostEqual(motion.ByX, 40.0)
        self.assertEqual(motion.ByY, 0.0)

    def test_reveal_starts_with_sentence_narration(self):
        audio = audio_effect(duration=3.0, delay=0.7)
        sequence = FakeSequence([audio])
        _, sequence, _ = self._process_sentence(sequence=sequence)
        reveal = sequence.effects[1]

        self.assertIs(sequence.effects[0], audio)
        self.assertEqual(
            reveal.Timing.TriggerType,
            self.processor.WITH_PREVIOUS,
        )
        self.assertEqual(
            reveal.Timing.TriggerDelayTime,
            self.settings.sentence_delay,
        )
        self.assertEqual(audio.Timing.TriggerDelayTime, 0.7)

    def test_configured_reveal_duration_is_used(self):
        settings = replace(self.settings, sentence_duration=0.77)
        planner = VisualAnimationPlanner(settings=settings)
        processor = VisualAnimationPresentationProcessor(
            planner=planner
        )
        spec = planner.plan_slide(
            FakeTemplateSlide(
                [FakeTemplateShape("PAST_SENTENCE")]
            ),
            2,
        )[0]
        _, sequence, _ = self._process_sentence(
            processor=processor,
            spec=spec,
        )

        self.assertEqual(sequence.effects[1].Timing.Duration, 0.77)

    def test_reveal_duration_does_not_exceed_audio_duration(self):
        sequence = FakeSequence([audio_effect(duration=0.8)])
        _, sequence, _ = self._process_sentence(
            slide_within_word=4,
            sequence=sequence,
        )

        self.assertEqual(sequence.effects[1].Timing.Duration, 0.8)

    def test_audio_metadata_and_order_remain_unchanged(self):
        audio_effects = [
            audio_effect(0.91, 0.5),
            audio_effect(3.47, 0.3),
        ]
        sequence = FakeSequence(audio_effects)
        slide = FakeSlide(
            shapes=[FakeShape("PAST_SENTENCE")],
            sequence=sequence,
        )
        before = [
            (
                effect.Timing.TriggerType,
                effect.Timing.TriggerDelayTime,
                effect.Timing.Duration,
            )
            for effect in audio_effects
        ]

        self.processor.process_slide(
            slide,
            [self._sentence_spec(2)],
            2,
            2,
        )

        after = [
            (
                effect.Timing.TriggerType,
                effect.Timing.TriggerDelayTime,
                effect.Timing.Duration,
            )
            for effect in audio_effects
        ]
        remaining_audio = [
            effect
            for effect in sequence.effects
            if effect.EffectType == 83
        ]

        self.assertEqual(before, after)
        self.assertEqual(remaining_audio, audio_effects)
        self.assertEqual(slide.SlideShowTransition.AdvanceTime, 7.25)

    def test_progress_shapes_remain_untouched(self):
        progress_track = FakeShape(
            "PROGRESS_TRACK",
            left=5,
            top=5,
            width=800,
            height=10,
        )
        progress_fill = FakeShape(
            "PROGRESS_FILL",
            left=5,
            top=5,
            width=200,
            height=10,
        )
        progress_before = [
            (shape.Left, shape.Top, shape.Width, shape.Height)
            for shape in (progress_track, progress_fill)
        ]
        _, sequence, _ = self._process_sentence(
            extra_shapes=[progress_track, progress_fill]
        )

        self.assertEqual(
            [
                (shape.Left, shape.Top, shape.Width, shape.Height)
                for shape in (progress_track, progress_fill)
            ],
            progress_before,
        )
        self.assertTrue(
            all(
                effect.Shape not in (progress_track, progress_fill)
                for effect in sequence.effects
                if effect.Shape is not None
            )
        )

    def test_mask_does_not_target_unrelated_element(self):
        unrelated = FakeShape(
            "DECORATIVE_LOGO",
            left=700,
            top=20,
            width=100,
            height=40,
        )
        unrelated_before = (
            unrelated.Left,
            unrelated.Top,
            unrelated.Width,
            unrelated.Height,
        )
        slide, sequence, sentence = self._process_sentence(
            extra_shapes=[unrelated]
        )
        mask = slide.Shapes.items[-1]

        self.assertIs(sequence.effects[1].Shape, mask)
        self.assertIsNot(sequence.effects[1].Shape, unrelated)
        self.assertEqual(
            (mask.Left, mask.Top, mask.Width, mask.Height),
            (sentence.Left, sentence.Top, sentence.Width, sentence.Height),
        )
        self.assertEqual(
            (
                unrelated.Left,
                unrelated.Top,
                unrelated.Width,
                unrelated.Height,
            ),
            unrelated_before,
        )

    def test_actual_reveal_mask_is_reported_once(self):
        output = io.StringIO()

        with redirect_stdout(output):
            self._process_sentence()

        self.assertEqual(output.getvalue().count("PAST_SENTENCE"), 1)
        self.assertIn("reveal mask (1.20s)", output.getvalue())

    def test_failed_mask_construction_is_clear_and_cleans_up(self):
        sequence = FakeSequence(
            [audio_effect(duration=3.0)],
            behavior_failure=True,
        )
        slide = FakeSlide(
            shapes=[FakeShape("PAST_SENTENCE")],
            sequence=sequence,
        )

        with self.assertRaisesRegex(
            AnimationTemplateError,
            "could not construct.*PAST_SENTENCE",
        ):
            self.processor.process_slide(
                slide,
                [self._sentence_spec(2)],
                2,
                2,
            )

        self.assertEqual(sequence.Count, 1)
        self.assertEqual(slide.Shapes.Count, 1)

    def test_repeated_static_content_is_not_animated_on_slides_2_to_4(self):
        forbidden = {
            "word",
            "pronunciation",
            "meaning",
            "translation",
            "verb_form",
            "image",
        }

        for slide_plan in self.template_plans[1:]:
            self.assertTrue(
                forbidden.isdisjoint(
                    {spec.semantic_element for spec in slide_plan}
                )
            )

    def test_fade_applies_only_to_first_slide_of_each_word(self):
        fade_slides = [
            slide_index
            for slide_index in range(1, 41)
            if self.planner.transition_spec(
                self.planner.slide_within_word(slide_index)
            ).entry_effect
            == self.settings.transition_id(
                self.settings.new_word_transition
            )
        ]

        self.assertEqual(fade_slides, list(range(1, 41, 4)))

    def test_slides_2_to_4_have_no_transition(self):
        for slide_within_word in (2, 3, 4):
            slide = FakeSlide()
            before_advance = (
                slide.SlideShowTransition.AdvanceOnTime,
                slide.SlideShowTransition.AdvanceTime,
            )
            self.processor.process_slide(
                slide,
                [],
                slide_within_word,
                slide_within_word,
            )

            self.assertEqual(
                slide.SlideShowTransition.EntryEffect,
                self.settings.transition_id(
                    self.settings.continuation_transition
                ),
            )
            self.assertEqual(
                (
                    slide.SlideShowTransition.AdvanceOnTime,
                    slide.SlideShowTransition.AdvanceTime,
                ),
                before_advance,
            )

    def test_sentence_semantic_name_must_select_exactly_one_target(self):
        duplicate_slide = FakeTemplateSlide(
            [
                FakeTemplateShape("PAST_SENTENCE"),
                FakeTemplateShape("PAST_SENTENCE"),
            ]
        )

        with self.assertRaisesRegex(
            AnimationTemplateError,
            "exactly one.*PAST_SENTENCE",
        ):
            self.planner.plan_slide(duplicate_slide, 2)

    def test_missing_required_sentence_shape_raises_template_error(self):
        with self.assertRaisesRegex(
            AnimationTemplateError,
            "PRESENT_SENTENCE",
        ):
            self.planner.plan_slide(FakeTemplateSlide(), 3)

    def test_missing_required_rendered_sentence_raises_template_error(self):
        with self.assertRaisesRegex(
            AnimationTemplateError,
            "FUTURE_SENTENCE",
        ):
            self.processor.process_slide(
                FakeSlide(
                    sequence=FakeSequence([audio_effect()])
                ),
                [self._sentence_spec(4)],
                4,
                4,
            )

    def test_missing_optional_intro_shape_does_not_damage_timeline(self):
        audio = audio_effect()
        sequence = FakeSequence([audio])
        slide = FakeSlide(sequence=sequence)
        spec = VisualAnimationSpec(
            shape_name="Optional Intro",
            semantic_element="meaning",
            effect_id=self.settings.effect_id(
                self.settings.meaning_effect
            ),
            duration=self.settings.meaning_duration,
        )

        self.processor.process_slide(slide, [spec], 1, 1)

        self.assertEqual(sequence.effects, [audio])

    def test_invalid_negative_duration_is_rejected(self):
        with self.assertRaisesRegex(
            AnimationConfigurationError,
            "ANIMATION_SENTENCE_DURATION.*negative",
        ):
            replace(self.settings, sentence_duration=-0.1)

    def test_invalid_mask_color_is_rejected(self):
        with self.assertRaisesRegex(
            AnimationConfigurationError,
            "ANIMATION_SENTENCE_MASK_COLOR",
        ):
            replace(self.settings, sentence_mask_color="white")

    def test_progress_names_remain_excluded_from_intro_planning(self):
        for shape_name in ("PROGRESS_TRACK", "PROGRESS_FILL"):
            self.assertIsNone(
                self.planner.plan_intro_shape(
                    shape_name,
                    "{{WORD}}",
                )
            )


if __name__ == "__main__":
    unittest.main()
