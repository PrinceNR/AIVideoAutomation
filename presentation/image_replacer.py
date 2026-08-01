# from pathlib import Path

# from pptx.enum.shapes import MSO_SHAPE_TYPE


# class ImageReplacer:

#     IMAGE_LEFT = 1223367
#     IMAGE_TOP = 728662

#     def replace_image(
#         self,
#         slide,
#         image_path: Path
#     ):

#         picture = self._find_picture(slide)

#         if picture is None:
#             return

#         left = picture.left
#         top = picture.top
#         width = picture.width
#         height = picture.height

#         # Remove old image
#         picture._element.getparent().remove(
#             picture._element
#         )

#         # Insert new image
#         slide.shapes.add_picture(
#             str(image_path),
#             left,
#             top,
#             width,
#             height
#         )



from pathlib import Path

from pptx.enum.shapes import MSO_SHAPE_TYPE


class ImageReplacer:

    IMAGE_LEFT = 1223367
    IMAGE_TOP = 728662

    def replace_image(
        self,
        slide,
        image_path: Path
    ):

        picture = self._find_picture(slide)

        if picture is None:
            return

        left = picture.left
        top = picture.top
        width = picture.width
        height = picture.height

        # Remove old image
        picture._element.getparent().remove(
            picture._element
        )

        # Insert new image
        slide.shapes.add_picture(
            str(image_path),
            left,
            top,
            width,
            height
        )

    def _find_picture(self, slide):

        for i, shape in enumerate(slide.shapes):

            if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:

                print(
                    f"Picture {i}:",
                    shape.name,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )

        return None
    # def _find_picture(self, slide):

    #     for shape in slide.shapes:

    #         if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
    #             continue

    #         if (
    #             shape.left == self.IMAGE_LEFT
    #             and
    #             shape.top == self.IMAGE_TOP
    #         ):
    #             return shape

    #     return None

    # def _find_picture(self, slide):

    #     for shape in slide.shapes:

    #         if shape.shape_type != MSO_SHAPE_TYPE.PICTURE:
    #             continue

    #         if (
    #             shape.left == self.IMAGE_LEFT
    #             and
    #             shape.top == self.IMAGE_TOP
    #         ):
    #             return shape

    #     return None